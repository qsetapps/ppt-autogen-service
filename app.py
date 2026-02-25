import io
import os
import datetime as dt
from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from openpyxl import load_workbook

app = Flask(__name__)

MONTHS_SV = [
    "Januari","Februari","Mars","April","Maj","Juni",
    "Juli","Augusti","September","Oktober","November","December"
]

def to_period_label(excel_value):
    if isinstance(excel_value, (dt.datetime, dt.date)):
        d = excel_value.date() if isinstance(excel_value, dt.datetime) else excel_value
        return f"{MONTHS_SV[d.month-1]} {d.year}"

    s = str(excel_value).strip()
    for fmt in ("%Y-%m-%d", "%d/%m/%y", "%m/%d/%y", "%d/%m/%Y", "%m/%d/%Y"):
        try:
            d = dt.datetime.strptime(s, fmt).date()
            return f"{MONTHS_SV[d.month-1]} {d.year}"
        except ValueError:
            pass
    return s  # fallback

def set_shape_text(slide, shape_id, new_text):
    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.has_text_frame:
            shape.text = str(new_text)
            return True
    return False

def set_labeled_value(slide, shape_id, label, value, keep_rest_lines=True):
    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.has_text_frame:
            raw = shape.text or ""
            normalized = raw.replace("\x0b", "\n")
            lines = normalized.split("\n")

            out = [label, str(value)]
            if keep_rest_lines and len(lines) > 2:
                out.extend(lines[2:])

            shape.text = "\n".join(out).replace("\n", "\x0b")
            return True
    return False

def update_ppt(excel_bytes: bytes, ppt_bytes: bytes) -> bytes:
    wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
    if "Auto" not in wb.sheetnames:
        raise ValueError("Saknar flik 'Auto' i Excel.")

    ws = wb["Auto"]

    # Slide 1 label
    period_label = to_period_label(ws["A1"].value)

    prs = Presentation(io.BytesIO(ppt_bytes))

    if len(prs.slides) < 2:
        raise ValueError("PPT måste ha minst 2 slides.")

    slide1 = prs.slides[0]
    slide2 = prs.slides[1]

    # Slide 1: shape_id=2 rubrik
    if not set_shape_text(slide1, 2, period_label):
        raise ValueError("Hittade inte shape_id=2 på slide 1.")

    # Läs budget-data från Auto!B4:F6
    oms, tg1, tg2, tg3, ebita = [ws[c].value for c in ["B4","C4","D4","E4","F4"]]
    oms_g, tg1_g, tg2_g, tg3_g, ebita_g = [ws[c].value for c in ["B5","C5","D5","E5","F5"]]
    oms_gsek, tg1_gsek, tg2_gsek, tg3_gsek, ebita_gsek = [ws[c].value for c in ["B6","C6","D6","E6","F6"]]

    # Slide 2 shape-ids (från din PPT)
    set_labeled_value(slide2, 3,  "Omsättning",   oms)
    set_labeled_value(slide2, 26, "TG 1",         tg1)
    set_labeled_value(slide2, 27, "TG 2",         tg2)
    set_labeled_value(slide2, 28, "TG 3/EBITDA",  tg3)
    set_labeled_value(slide2, 29, "EBITA",        ebita)

    set_labeled_value(slide2, 10, "Tillväxt",     oms_g, keep_rest_lines=False)
    set_labeled_value(slide2, 13, "Tillväxt",     tg1_g, keep_rest_lines=False)
    set_labeled_value(slide2, 17, "Tillväxt",     tg2_g, keep_rest_lines=False)
    set_labeled_value(slide2, 32, "Tillväxt",     tg3_g, keep_rest_lines=False)
    set_labeled_value(slide2, 35, "Tillväxt",     ebita_g, keep_rest_lines=False)

    set_labeled_value(slide2, 11, "Tillväxt Sek", oms_gsek, keep_rest_lines=False)
    set_labeled_value(slide2, 15, "Tillväxt Sek", tg1_gsek, keep_rest_lines=False)
    set_labeled_value(slide2, 31, "Tillväxt Sek", tg2_gsek, keep_rest_lines=False)
    set_labeled_value(slide2, 33, "Tillväxt Sek", tg3_gsek, keep_rest_lines=False)
    set_labeled_value(slide2, 36, "Tillväxt Sek", ebita_gsek, keep_rest_lines=False)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def require_token():
    expected = os.environ.get("API_TOKEN")
    if not expected:
        return  # allow if not set (dev)
    got = request.headers.get("X-API-Token", "")
    if got != expected:
        raise PermissionError("Unauthorized")

@app.route("/health", methods=["GET"])
def health():
    return {"ok": True}

@app.route("/update", methods=["POST"])
def update():
    try:
        require_token()

        # Vi förväntar oss multipart/form-data med fälten:
        # - excel: .xlsx
        # - ppt: .pptx
        if "excel" not in request.files or "ppt" not in request.files:
            return jsonify({
                "error": "Skicka multipart/form-data med filer i fälten 'excel' och 'ppt'."
            }), 400

        excel_bytes = request.files["excel"].read()
        ppt_bytes = request.files["ppt"].read()

        out_bytes = update_ppt(excel_bytes, ppt_bytes)

        return send_file(
            io.BytesIO(out_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name="Ledningsmote_autogen.pptx",
        )

    except PermissionError as e:
        return jsonify({"error": str(e)}), 401
    except Exception as e:
        return jsonify({"error": str(e)}), 500
